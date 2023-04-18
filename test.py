import os
import pptx
from pptx.util import Inches

current_directory = os.getcwd()
replacement_animals = os.listdir(current_directory + '/replacement_images')
prs = pptx.Presentation('original.pptx')

image_to_replace = pptx.parts.image.Image.from_file('lion.png')

# Find the slide you want to modify
slide = prs.slides[0]

replacement_counter = 0

shape = slide.shapes[1]
# Iterate though the shapes looking for dogs
for shape in slide.shapes:
    if shape.name == 'Lion': # manually named all the lion images to be "Lion" in the pptx file
        if shape.shape_type == 13:  # 13 means it is a picture

            # print(shape.image)
            # print(shape)
            # print(f'replacement_imagesdes/{replacement_animalsdes[replacement_counter]}')
            
            # print(Inches(shape.left))
            # print(Inches(shape.top))
            # print(Inches(shape.width))
            # print(Inches(shape.height))

            # new_shape = slide.shapes.add_picture(
            #     f'replacement_imagesdes/{replacement_animalsdes[replacement_counter]}',
            #     Inches(shape.left),
            #     Inches(shape.top),
            #     Inches(shape.width),
            #     Inches(shape.height),
            # )
        
            # print(new_shape._element)
            # old_pic = shape._element
            # new_pic = new_shape._element
            # old_pic.addnext(new_pic)
            # old_pic.getparent().remove(old_pic)

            print(f'replacement_images/{replacement_animals[replacement_counter]}')
            new_pptx_img = pptx.parts.image.Image.from_file(f'replacement_images/{replacement_animals[replacement_counter]}')

            # shape._element.blip_rId = 'rId' + str(replacement_counter + 300)

            slide_part, rId = shape.part, shape._element.blip_rId

            print('rID', shape._element.blip_rId)
            image_part = slide_part.related_part(rId)

            if image_to_replace.blob == image_part.blob:
                image_part.blob = new_pptx_img._blob

            # image_part.blob = new_pptx_img._blob

            replacement_counter += 1

            # prs.save('edited.pptx')


# Save the modified PowerPoint file
prs.save('edited.pptx')